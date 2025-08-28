#!/usr/bin/env python3
"""
Create PowerPoint Presentation for Math Solver App
Enhanced version with proper framing and organization
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR

def add_formatted_text(text_frame, content, font_size=18, is_title=False):
    """Add properly formatted text to a text frame"""
    text_frame.clear()
    
    if isinstance(content, str):
        # Split into lines and process
        lines = content.split('\n')
    else:
        lines = content
    
    for i, line in enumerate(lines):
        if line.strip():  # Skip empty lines
            if i > 0:  # Add paragraph for new lines
                p = text_frame.add_paragraph()
            else:
                p = text_frame.paragraphs[0]
            
            p.text = line.strip()
            p.font.size = Pt(font_size)
            
            # Set bullet style for non-title content
            if not is_title and line.startswith(('•', '-', '✅', '📱', '🤖', '⚡')):
                p.level = 0  # Bullet level
            elif not is_title and line.startswith(('    •', '    -')):
                p.level = 1  # Sub-bullet level

def create_slide_with_content(prs, layout_type, title_text, content_text, title_size=32, content_size=18):
    """Create a well-formatted slide with proper spacing"""
    slide_layout = prs.slide_layouts[layout_type]
    slide = prs.slides.add_slide(slide_layout)
    
    # Configure title
    if slide.shapes.title:
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.size = Pt(title_size)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 122, 255)
        title.text_frame.paragraphs[0].font.bold = True
        
        # Ensure title is properly positioned
        title.top = Inches(0.5)
        title.left = Inches(0.5)
        title.width = Inches(12.33)
        title.height = Inches(1.2)
    
    # Configure content
    if len(slide.placeholders) > 1:
        content = slide.placeholders[1]
        add_formatted_text(content.text_frame, content_text, content_size)
        
        # Ensure content is properly positioned and sized
        content.top = Inches(1.8)
        content.left = Inches(0.5)
        content.width = Inches(12.33)
        content.height = Inches(5.2)
    
    return slide

def create_math_solver_presentation():
    """Create a comprehensive PowerPoint presentation for the Math Solver App"""
    
    # Create presentation object
    prs = Presentation()
    
    # Set slide dimensions (16:9 aspect ratio)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Define colors
    blue_color = RGBColor(0, 122, 255)  # Primary blue
    dark_blue = RGBColor(0, 64, 133)    # Dark blue
    green_color = RGBColor(40, 167, 69)  # Success green
    orange_color = RGBColor(255, 193, 7) # Warning orange
    
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    # Configure title
    title.text = "Math Solver 🧮"
    title.text_frame.paragraphs[0].font.size = Pt(48)
    title.text_frame.paragraphs[0].font.color.rgb = blue_color
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Position title properly
    title.top = Inches(1.5)
    title.left = Inches(1)
    title.width = Inches(11.33)
    title.height = Inches(1.5)
    
    # Configure subtitle
    subtitle.text = "AI-Powered Mathematical Problem Solver\n\nFrom Image to Solution in Seconds\n\nTechnology Stack: React Native + FastAPI + OpenAI GPT-4o-mini"
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Position subtitle properly
    subtitle.top = Inches(3.5)
    subtitle.left = Inches(1)
    subtitle.width = Inches(11.33)
    subtitle.height = Inches(3)
    
    # Slide 2: Problem Statement
    create_slide_with_content(prs, 1, "The Challenge 🤔", 
"""• Students struggle with complex mathematical problems
• Time-consuming manual calculations and verification
• Limited access to step-by-step guidance
• Need for instant mathematical assistance
• Gap between problem recognition and solution methodology""")
    
    # Slide 3: Our Solution
    create_slide_with_content(prs, 1, "Math Solver App ✨", 
""""Snap a photo, get instant step-by-step solutions"

• 📸 Image Recognition: Take photos of handwritten or printed math problems
• 🤖 AI-Powered: Uses advanced OpenAI GPT-4o-mini vision model
• 📋 Step-by-Step: Detailed solution breakdown with explanations
• ⚡ Real-time: Results in under 90 seconds
• 📱 Cross-Platform: Works on iOS, Android, and Web""")
    
    # Slide 4: Key Features
    create_slide_with_content(prs, 1, "What Makes It Special 🌟", 
"""Core Features:
• 🔍 Smart Image Processing: Recognizes equations from photos
• 📊 Multiple Problem Types: Algebra, geometry, calculus support
• 📝 Detailed Explanations: Every step explained clearly
• ✅ Solution Verification: Confirms answer accuracy
• 🎯 User-Friendly Interface: Clean, intuitive design

Technical Features:
• ⚡ Fast Processing: 90-second timeout with progress tracking
• 🔒 Secure: API keys protected, no data leakage
• 🌐 Cross-Platform: React Native for universal compatibility""")
    
    # Slide 5: Technology Architecture
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "How It Works 🏗️"
    content.text = """Architecture Flow:
📱 Frontend (React Native)
    ↓ Image Upload
🔗 API Gateway (FastAPI)
    ↓ Image Processing
🤖 AI Processing (OpenAI GPT-4o-mini)
    ↓ JSON Response
📊 Solution Display (Structured)

Tech Stack:
• Frontend: React Native + Expo
• Backend: Python FastAPI
• AI Engine: OpenAI GPT-4o-mini Vision Model
• Infrastructure: RESTful API architecture"""
    
    # Slide 6: User Journey
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "From Problem to Solution 🎯"
    content.text = """Step 1: Capture 📸
• Open the app
• Take photo or select from gallery
• Auto-crop and optimize image

Step 2: Process 🔄
• Upload to secure backend
• AI analyzes mathematical content
• Extract equations and context

Step 3: Solve 🧮
• GPT-4o-mini processes the problem
• Generates step-by-step solution
• Validates and formats response

Step 4: Learn 📚
• View detailed solution steps
• Understand methodology
• Verify final answer"""
    
    # Slide 7: Sample Problem Walkthrough
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Example: System of Equations 📊"
    content.text = """Input Image:
x + y = 31
x - y = 18

AI-Generated Solution:
1. Step 1: Add both equations to eliminate y
   • (x + y) + (x - y) = 31 + 18
   • 2x = 49

2. Step 2: Solve for x
   • x = 49 ÷ 2 = 24.5

3. Step 3: Substitute back to find y
   • 24.5 + y = 31
   • y = 6.5

Verification: ✅
• 24.5 + 6.5 = 31 ✓
• 24.5 - 6.5 = 18 ✓"""
    
    # Slide 8: Technical Highlights
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Engineering Excellence 🔧"
    content.text = """Performance Optimizations:
⚡ Single Request Queue: Prevents API overload
⏱️ 90-Second Timeout: Handles complex problems
🔄 Automatic Retry: Smart fallback mechanisms
📱 Smart URL Detection: Cross-platform connectivity

Security Features:
🔒 Environment Variables: API keys secured
🛡️ Input Validation: File type and size checks
🚫 No Data Storage: Privacy-first approach
📋 Comprehensive Logging: Error tracking and debugging"""
    
    # Slide 9: Market Opportunity
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Target Audience 🎯"
    content.text = """Primary Users:
📚 Students: High school and college mathematics
👨‍🏫 Educators: Teaching aids and verification tools
👨‍💼 Professionals: Engineering and scientific calculations
👨‍👩‍👧‍👦 Parents: Helping children with homework

Market Size:
📈 Global EdTech Market: $89.49 billion (2020)
📱 Mobile Learning: 37% annual growth rate
🤖 AI in Education: $25.7 billion by 2030"""
    
    # Slide 10: Competitive Advantage
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "What Sets Us Apart 🏆"
    content.text = """vs. Traditional Calculators:
✅ Image Recognition vs Manual Input
✅ Step-by-Step Explanations vs Just Answers
✅ Complex Problem Handling vs Basic Operations

vs. Existing Apps:
✅ Advanced AI (GPT-4o-mini) vs Basic OCR
✅ Cross-Platform Native vs Web-Only
✅ Detailed Verification vs Simple Results
✅ Modern Architecture vs Legacy Systems"""
    
    # Slide 11: Future Roadmap
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "What's Next 🔮"
    content.text = """Short Term (3-6 months):
📈 Enhanced AI Models: Support for more problem types
🌍 Multi-Language: Support for different mathematical notations
💾 History Feature: Save and review past solutions
📊 Analytics Dashboard: Usage insights and improvements

Long Term (6-12 months):
🎓 Educational Content: Integrated learning modules
👥 Collaboration Features: Share and discuss solutions
🏢 Enterprise Version: Bulk processing for institutions
🔗 API Marketplace: Third-party integrations"""
    
    # Slide 12: Business Model
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Monetization Strategy 💰"
    content.text = """Freemium Model:
📱 Free Tier: 10 problems per day
⭐ Premium Tier: Unlimited problems + priority processing
🎓 Educational Tier: Discounted rates for schools

Revenue Streams:
💳 Subscription Fees: Monthly/Annual premium plans
🏢 Enterprise Licenses: Institutional partnerships
📚 Educational Content: Premium learning materials
🔗 API Access: Developer integrations"""
    
    # Slide 13: Success Metrics
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Measuring Impact 📈"
    content.text = """User Engagement:
👥 Daily Active Users: Target 10K+ within 6 months
🔄 Problem Solving Rate: >95% success rate
⏱️ Average Response Time: <30 seconds
⭐ User Satisfaction: 4.5+ stars

Technical Performance:
🚀 API Uptime: 99.9% availability
⚡ Response Time: <90 seconds guarantee
🔒 Security Incidents: Zero tolerance
📊 Error Rate: <1% failure rate"""
    
    # Slide 14: Call to Action
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Next Steps 🎯"
    content.text = """For Investors:
💰 Funding Opportunity: Scalable EdTech solution
📈 Market Potential: Huge addressable market
🏆 Competitive Edge: Advanced AI integration

For Users:
📱 Try the App: Download and test today
📝 Feedback Welcome: Help us improve
📢 Spread the Word: Share with students and educators

For Developers:
🔗 API Access: Integration opportunities
👥 Collaboration: Open source contributions
🚀 Innovation: Building the future of education"""
    
    # Slide 15: Thank You
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Thank You! Questions & Discussion 🙋‍♂️"
    content.text = """Contact Information:
📧 Email: [your.email@domain.com]
💼 LinkedIn: [Your LinkedIn Profile]
🐱 GitHub: [Your GitHub Repository]
🌐 Demo: [Live Demo URL]

Resources:
📖 Documentation: Comprehensive setup guides
🎥 Video Demo: Full walkthrough available
📊 Technical Details: Architecture diagrams
🔗 Source Code: Available on GitHub"""
    
    # Save the presentation
    prs.save('Math_Solver_Presentation.pptx')
    print("✅ PowerPoint presentation created successfully!")
    print("📁 File saved as: Math_Solver_Presentation.pptx")
    
if __name__ == "__main__":
    try:
        create_math_solver_presentation()
    except ImportError as e:
        print("❌ Error: python-pptx library not found!")
        print("📦 Please install it using: pip install python-pptx")
        print(f"   Error details: {e}")
    except Exception as e:
        print(f"❌ Error creating presentation: {e}")
