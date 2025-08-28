#!/usr/bin/env python3
"""
Create Refined PowerPoint Presentation for Math Solver App
Focused on key points with minimal text for better presentation impact
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def add_formatted_text(text_frame, content, font_size=22):
    """Add properly formatted text to a text frame"""
    text_frame.clear()
    lines = content.split('\n') if isinstance(content, str) else content
    
    for i, line in enumerate(lines):
        if line.strip():
            if i > 0:
                p = text_frame.add_paragraph()
            else:
                p = text_frame.paragraphs[0]
            
            p.text = line.strip()
            p.font.size = Pt(font_size)
            p.font.name = 'Calibri'
            
            # Set bullet style for bullet points
            if line.strip().startswith(('•', '📸', '🤖', '⚡')):
                p.level = 0

def create_slide_with_content(prs, title_text, content_text, layout_type=1):
    """Create a well-formatted slide with proper spacing and framing"""
    slide_layout = prs.slide_layouts[layout_type]
    slide = prs.slides.add_slide(slide_layout)
    
    # Configure title
    if slide.shapes.title:
        title = slide.shapes.title
        title.text = title_text
        
        # Title formatting
        title_para = title.text_frame.paragraphs[0]
        title_para.font.size = Pt(40)
        title_para.font.color.rgb = RGBColor(0, 122, 255)
        title_para.font.bold = True
        title_para.font.name = 'Calibri'
        
        # Position title with proper margins
        title.top = Inches(0.3)
        title.left = Inches(0.5)
        title.width = Inches(12.33)
        title.height = Inches(1.0)
    
    # Configure content
    if len(slide.placeholders) > 1 and content_text:
        content = slide.placeholders[1]
        add_formatted_text(content.text_frame, content_text, 24)
        
        # Position content with proper margins
        content.top = Inches(1.5)
        content.left = Inches(0.5)
        content.width = Inches(12.33)
        content.height = Inches(5.7)
        
        # Set content text properties
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.name = 'Calibri'
            paragraph.space_after = Pt(12)  # Add spacing between points
    
    return slide

def create_math_solver_presentation():
    """Create refined PowerPoint presentation with key points only"""
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    # Title
    title.text = "Math Solver 🧮"
    title_para = title.text_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.color.rgb = RGBColor(0, 122, 255)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    
    title.top = Inches(1.8)
    title.left = Inches(0.5)
    title.width = Inches(12.33)
    title.height = Inches(1.5)
    
    # Subtitle
    subtitle.text = "AI-Powered Mathematical Problem Solver\n\nFrom Image to Solution in Seconds"
    sub_para = subtitle.text_frame.paragraphs[0]
    sub_para.font.size = Pt(28)
    sub_para.alignment = PP_ALIGN.CENTER
    
    subtitle.top = Inches(3.8)
    subtitle.left = Inches(0.5)
    subtitle.width = Inches(12.33)
    subtitle.height = Inches(2.5)
    
    # Slide 2: The Problem
    create_slide_with_content(prs, "The Problem 🤔", 
"""• Students struggle with complex math problems
• No instant step-by-step guidance available
• Time-consuming manual calculations""")
    
    # Slide 3: Our Solution
    create_slide_with_content(prs, "Our Solution ✨", 
"""📸 Snap a photo of any math problem

🤖 AI analyzes and solves instantly

📋 Get detailed step-by-step solutions""")
    
    # Slide 4: How It Works
    create_slide_with_content(prs, "How It Works 🔄", 
"""1️⃣ Take a photo of your math problem

2️⃣ AI processes the image in <90 seconds

3️⃣ Receive step-by-step solution with verification""")
    
    # Slide 5: Key Features
    create_slide_with_content(prs, "Key Features 🌟", 
"""⚡ Fast: Results in under 90 seconds

🎯 Accurate: Advanced AI with verification

📱 Easy: Simple photo-to-solution interface""")
    
    # Slide 6: Technology Stack
    create_slide_with_content(prs, "Technology Stack 🏗️", 
"""📱 React Native Frontend

🐍 FastAPI Backend

🤖 OpenAI GPT-4o-mini Vision Model""")
    
    # Slide 7: Live Demo Example
    create_slide_with_content(prs, "Example: System of Equations 📊", 
"""Input: x + y = 31, x - y = 18

Solution Steps:
• Add equations: 2x = 49
• Solve: x = 24.5, y = 6.5
• Verify: ✅ Both equations satisfied""")
    
    # Slide 8: Market Opportunity
    create_slide_with_content(prs, "Market Opportunity 🎯", 
"""📚 Students: 1.6 billion worldwide

👨‍🏫 Educators: Need better teaching tools

📈 EdTech Market: $89.49 billion globally""")
    
    # Slide 9: Business Model
    create_slide_with_content(prs, "Business Model 💰", 
"""📱 Freemium: 10 problems/day free

⭐ Premium: Unlimited + priority processing

🏢 Enterprise: Educational institutions""")
    
    # Slide 10: Competitive Advantage
    create_slide_with_content(prs, "Why We Win 🏆", 
"""🔍 Advanced AI vs basic calculators

📝 Step-by-step explanations vs just answers

📸 Image recognition vs manual input""")
    
    # Slide 11: Roadmap
    create_slide_with_content(prs, "What's Next 🔮", 
"""🌍 Multi-language support

📊 Advanced problem types

🎓 Educational content integration""")
    
    # Slide 12: Call to Action
    create_slide_with_content(prs, "Get Involved 🚀", 
"""💰 Investors: Join the EdTech revolution

📱 Users: Try the app today

👥 Developers: Contribute on GitHub

🐱 Repository: github.com/lskenkf/math_solver""")
    
    # Slide 13: Thank You
    create_slide_with_content(prs, "Thank You! 🙏", 
"""Questions & Discussion

📧 Contact: [your.email@domain.com]

🐱 GitHub: github.com/lskenkf/math_solver

🚀 Let's revolutionize math education together!""")
    
    # Save the presentation
    prs.save('Math_Solver_Refined_Presentation.pptx')
    print("✅ Refined PowerPoint presentation created successfully!")
    print("📁 File saved as: Math_Solver_Refined_Presentation.pptx")
    print("🎯 Features: Key points only, larger fonts, better presentation flow")

if __name__ == "__main__":
    try:
        create_math_solver_presentation()
    except ImportError as e:
        print("❌ Error: python-pptx library not found!")
        print("📦 Please install it using: pip install python-pptx")
        print(f"   Error details: {e}")
    except Exception as e:
        print(f"❌ Error creating presentation: {e}")
